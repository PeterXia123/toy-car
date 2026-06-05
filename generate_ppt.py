# -*- coding: utf-8 -*-
"""Generate a presentation showcasing the MVA Validator system."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

DARK = RGBColor(0x2C, 0x3E, 0x50)
BLUE = RGBColor(0x34, 0x98, 0xDB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
GREEN = RGBColor(0x27, 0xAE, 0x60)
GREY = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_GREY = RGBColor(0x95, 0xA5, 0xA6)
ACCENT = RGBColor(0x29, 0x80, 0xB9)


def _set_bg(slide, color=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_bar(slide, title):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.6)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK)

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.85), Inches(1.2), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(8.4), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = LIGHT_GREY
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(16)


def add_bullet_slide(prs, title, items):
    """items: list of (bold_key, description) tuples."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _add_title_bar(slide, title)

    y_start = Inches(1.5)
    spacing = Inches(0.95)

    for i, (key, desc) in enumerate(items):
        y = y_start + i * spacing

        # Accent dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.7), int(y + Pt(6)), Pt(10), Pt(10))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLUE
        dot.line.fill.background()

        # Text box with bold key + description
        txBox = slide.shapes.add_textbox(Inches(1.1), int(y), Inches(8.2), Inches(0.85))
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        run_key = p.add_run()
        run_key.text = key
        run_key.font.size = Pt(18)
        run_key.font.bold = True
        run_key.font.color.rgb = DARK

        p2 = tf.add_paragraph()
        run_desc = p2.add_run()
        run_desc.text = desc
        run_desc.font.size = Pt(14)
        run_desc.font.color.rgb = GREY
        p2.space_before = Pt(2)


def add_three_col_slide(prs, title, columns):
    """columns: list of (heading, items_list) for 3 columns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _add_title_bar(slide, title)

    col_w = Inches(2.8)
    gap = Inches(0.2)
    start_x = Inches(0.5)

    for i, (heading, items) in enumerate(columns):
        x = start_x + i * (col_w + gap)
        y = Inches(1.5)

        # Column card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(col_w), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = RGBColor(0xDE, 0xE2, 0xE6)
        card.line.width = Pt(1)

        # Heading
        hdr = slide.shapes.add_textbox(int(x + Inches(0.2)), int(y + Inches(0.2)),
                                       int(col_w - Inches(0.4)), Inches(0.5))
        tf = hdr.text_frame
        p = tf.paragraphs[0]
        p.text = heading
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.alignment = PP_ALIGN.CENTER

        # Separator
        sep = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(x + Inches(0.4)), int(y + Inches(0.8)),
            int(col_w - Inches(0.8)), Pt(2))
        sep.fill.solid()
        sep.fill.fore_color.rgb = BLUE
        sep.line.fill.background()

        # Items
        body = slide.shapes.add_textbox(int(x + Inches(0.25)), int(y + Inches(1.0)),
                                        int(col_w - Inches(0.5)), Inches(3.5))
        btf = body.text_frame
        btf.word_wrap = True
        for j, item in enumerate(items):
            p = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
            p.text = item
            p.font.size = Pt(13)
            p.font.color.rgb = DARK
            p.space_before = Pt(8)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    add_title_slide(
        prs,
        "MVA Validator",
        "Automated Performance Data Validation for IFRS 9 Models"
    )

    # --- Slide 2: Starting Point ---
    add_bullet_slide(prs, "Starting Point", [
        ("200+ models validated per year",
         "Every model validation at IB starts with the same foundation: performance data"),
        ("One consolidated dataset per portfolio",
         "All portfolios share a common perf data structure in parquet format"),
        ("Data issues are the #1 source of rework",
         "Manual review of millions of observations is slow, inconsistent, and error-prone"),
    ])

    # --- Slide 3: One Command ---
    add_bullet_slide(prs, "One Command, Full Output", [
        ("One parquet file + one config",
         "Point the tool at your data file and a YAML config - that's all it needs"),
        ("Auto-discover issues with materiality",
         "Generates issue log with affected account %, impact level, and account-level examples"),
        ("Ready for EST or validation guidance",
         "Output can go directly into the EST package, or guide targeted deep-dive analysis"),
    ])

    # --- Slide 4: Three Categories (3-column layout) ---
    add_three_col_slide(prs, "Three Categories of Checks", [
        ("Variable", [
            "Data type validation",
            "Missing & negative values",
            "Outlier detection",
            "Indicator logic consistency",
            "(default, closed, charge-off)",
        ]),
        ("Trend", [
            "Default rate trends",
            "Score distribution drift",
            "Attrition spike detection",
            "Recovery trend monitoring",
            "(data migration signals)",
        ]),
        ("Connection", [
            "Score vs default alignment",
            "Term vs maturity consistency",
            "Balance vs recovery logic",
            "DPD vs default status",
            "(cross-variable checks)",
        ]),
    ])

    # --- Slide 5: IFRS 9 Parameter Tracking ---
    add_bullet_slide(prs, "Mapped to IFRS 9 Parameters", [
        ("Tagged per parameter",
         "Every check is linked to PD, LGD, EAD, ERL, DF, or SICR"),
        ("Filterable and trackable",
         "Review findings by parameter to support targeted validation per model component"),
        ("Tests designed independently per parameter",
         "Each IFRS 9 parameter has its own set of checks that can be run and tracked separately"),
    ])

    # --- Slide 6: Extensibility ---
    add_bullet_slide(prs, "Built to Extend", [
        ("Config-driven, no code changes",
         "New portfolio = new YAML file. Supports term (MG, AL, PL) and revolving (CC)"),
        ("Modular architecture",
         "Easy to add new checks per category without touching existing logic"),
        ("Future roadmap",
         "Plug into Commercial Banking portfolios, extend to AIRB, delta reporting across cycles"),
    ])

    # --- Slide 7: Thank You ---
    add_title_slide(prs, "Thank You", "Questions?")

    out_path = "output/MVA_Validator_Overview.pptx"
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    print(f"Saved to {out_path}")


if __name__ == "__main__":
    main()
